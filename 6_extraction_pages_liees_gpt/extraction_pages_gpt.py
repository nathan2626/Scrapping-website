import json
import os
import time
import re
import pymupdf as fitz  # PyMuPDF
import docx
import pptx
from openai import OpenAI

# ===========================
# 📂 Chemins et Paramètres
# ===========================
INPUT_DIR = "5_scraping_pages_liees"
OUTPUT_DIR = "6_extraction_pages_gpt"
LOG_FILE = os.path.join(OUTPUT_DIR, "gpt_pages_errors.log")
MAX_CONTENT = 10000

# Mots-clés pour filtrer uniquement les pages pertinentes
MOTS_CLES = [
    "congre", "congres", "congress", "congresse", "congresses",
    "2025", "2026", "2027", "event", "events", "planning"
]

# ===========================
# 💬 Initialisation GPT
# ===========================
with open('key.txt', 'r') as file:
    secret_key = file.read().strip()

client = OpenAI(api_key=secret_key)

# ===========================
# 📝 Gestion des logs
# ===========================
def log_erreur(message):
    """Enregistre une erreur dans le fichier de log."""
    with open(LOG_FILE, "a", encoding="utf-8") as f:
        f.write(f"{message}\n")
    print(f"❌ {message}")

# ===========================
# 📂 Filtrage des pages par mots-clés
# ===========================
def contient_mot_cle(texte):
    """Vérifie si une chaîne contient un mot-clé pertinent."""
    if not texte:
        return False
    texte = texte.lower()
    return any(mot in texte for mot in MOTS_CLES)

# ===========================
# 📖 Lecture des fichiers téléchargés
# ===========================
def lire_fichier(filepath):
    """Lit le contenu d’un fichier selon son type."""
    if filepath.endswith(".pdf"):
        return lire_pdf(filepath)
    elif filepath.endswith(".docx"):
        return lire_docx(filepath)
    elif filepath.endswith(".pptx"):
        return lire_pptx(filepath)
    return ""

def lire_pdf(filepath):
    """Lit le contenu d’un fichier PDF."""
    try:
        doc = fitz.open(filepath)
        return "\n".join([page.get_text() for page in doc])
    except Exception as e:
        log_erreur(f"Erreur lecture PDF {filepath}: {e}")
        return ""

def lire_docx(filepath):
    """Lit le contenu d’un fichier DOCX."""
    try:
        doc = docx.Document(filepath)
        return "\n".join([p.text for p in doc.paragraphs])
    except Exception as e:
        log_erreur(f"Erreur lecture DOCX {filepath}: {e}")
        return ""

def lire_pptx(filepath):
    """Lit le contenu d’un fichier PPTX."""
    try:
        presentation = pptx.Presentation(filepath)
        texte = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texte += shape.text + "\n"
        return texte
    except Exception as e:
        log_erreur(f"Erreur lecture PPTX {filepath}: {e}")
        return ""

# ===========================
# 💬 Analyse du contenu avec GPT
# ===========================
def analyser_contenu_avec_gpt(url, contenu):
    """Envoie le contenu à GPT pour analyse et retourne un JSON structuré."""
    prompt = f"""
    Voici un extrait d'une page web issue de {url}. Analyse-le et retourne **uniquement** un JSON valide en français, contenant :

    {{
        "contacts": [
            {{
                "type": "email",
                "valeur": "email@example.com",
                "propriétaire": {{
                    "nom": "Dr Jean Dupont",
                    "fonction": "Conférencier",
                    "affiliation": "Université de Paris"
                }}
            }},
            {{
                "type": "téléphone",
                "valeur": "+33 1 23 45 67 89",
                "propriétaire": {{
                    "nom": "Université de Paris",
                    "type": "Organisation"
                }}
            }},
            {{
                "type": "site_web",
                "valeur": "https://congres2025.fr",
                "propriétaire": {{
                    "nom": "Organisateur Congrès 2025"
                }}
            }}
        ],
        "personnes": [
            {{
                "nom": "Dr Jean Dupont",
                "titre": "Professeur",
                "affiliation": "Université de Paris",
                "contact_email": "jean.dupont@example.com",
                "contact_téléphone": "+33 1 23 45 67 89",
                "rôle": "Conférencier principal",
                "publications": [
                    "Étude sur l'IA médicale (2023)",
                    "Article dans ScienceDirect (2024)"
                ]
            }},
            {{
                "nom": "Marie Curie",
                "titre": "Chercheuse",
                "affiliation": "Institut Curie",
                "rôle": "Intervenante"
            }}
        ],
        "organisateurs": [
            {{
                "nom": "Société Française de Médecine",
                "type": "Association",
                "site_web": "https://sfmed.fr",
                "email_contact": "contact@sfmed.fr"
            }}
        ],
        "lieux": [
            {{
                "nom": "Centre des Congrès de Paris",
                "adresse": "Place de la République, 75003 Paris",
                "pays": "France"
            }}
        ],
        "dates": [
            {{
                "date_début": "2025-03-10",
                "date_fin": "2025-03-12",
                "événement": "Congrès International de Médecine",
                "personnes_associées": ["Dr Jean Dupont", "Marie Curie"],
                "programme": [
                    {{"heure": "09:00", "session": "Ouverture"}},
                    {{"heure": "14:00", "session": "Table ronde IA & Santé"}}
                ]
            }}
        ],
        "documents": [
            {{
                "type": "programme",
                "url": "https://congres2025.fr/programme.pdf"
            }},
            {{
                "type": "actes",
                "url": "https://congres2025.fr/actes.pdf"
            }}
        ],
        "réseaux_sociaux": [
            {{
                "plateforme": "Twitter",
                "lien": "https://twitter.com/congres2025"
            }},
            {{
                "plateforme": "LinkedIn",
                "lien": "https://linkedin.com/company/congres2025"
            }}
        ],
        "résumé": "Le Congrès International de Médecine 2025 se tiendra à Paris et abordera les avancées médicales avec des conférences d’experts et des tables rondes sur l’IA en santé."
    }}

    Extrait :
    {contenu[:MAX_CONTENT]}
    """

    try:
        response = client.chat.completions.create(
            model="gpt-4",
            messages=[{"role": "user", "content": prompt}]
        )
        resultat = response.choices[0].message.content.strip()

        if resultat.startswith("{") and resultat.endswith("}"):
            return json.loads(resultat)
        else:
            log_erreur(f"Erreur JSON GPT sur {url} : {resultat[:500]}")
            return {"erreur": "Réponse non JSON", "contenu_brut": resultat}

    except Exception as e:
        log_erreur(f"Erreur GPT sur {url}: {e}")
        return {"erreur": str(e)}

# ===========================
# 💾 Fonction principale
# ===========================
def analyser_toutes_pages():
    """Parcourt toutes les pages et fichiers téléchargés et lance l'analyse GPT."""
    for domaine in os.listdir(INPUT_DIR):
        domaine_path = os.path.join(INPUT_DIR, domaine)
        if not os.path.isdir(domaine_path):
            continue

        print(f"\n🚀 Analyse GPT pour le domaine : {domaine}")
        analyses = []

        # ===========================
        # 📖 Analyse des Pages HTML
        # ===========================
        pages_dir = os.path.join(domaine_path, "pages_html")
        if os.path.exists(pages_dir):
            for page_file in os.listdir(pages_dir):
                with open(os.path.join(pages_dir, page_file), "r", encoding="utf-8") as f:
                    contenu = f.read()

                # Filtrage par mots-clés
                if not contient_mot_cle(contenu) and not contient_mot_cle(page_file):
                    print(f"🚫 Ignoré (page non pertinente) : {page_file}")
                    continue

                print(f"💬 Analyse GPT de la page : {page_file}")
                resultat = analyser_contenu_avec_gpt(page_file, contenu)
                analyses.append({
                    "type": "page_html",
                    "nom": page_file,
                    "contenu": contenu[:500],
                    "analyse": resultat
                })

        # ===========================
        # 📂 Analyse des Fichiers Téléchargés
        # ===========================
        fichiers_dir = os.path.join(domaine_path, "fichiers")
        if os.path.exists(fichiers_dir):
            for fichier in os.listdir(fichiers_dir):
                # Exclure les images
                if re.search(r"\.(jpg|jpeg|png|gif|svg|webp)$", fichier, re.IGNORECASE):
                    print(f"🚫 Ignoré (image) : {fichier}")
                    continue

                contenu = lire_fichier(os.path.join(fichiers_dir, fichier))

                # Filtrage par mots-clés
                if not contient_mot_cle(contenu) and not contient_mot_cle(fichier):
                    print(f"🚫 Ignoré (fichier non pertinent) : {fichier}")
                    continue

                print(f"💬 Analyse GPT du fichier : {fichier}")
                resultat = analyser_contenu_avec_gpt(fichier, contenu)
                analyses.append({
                    "type": "fichier",
                    "nom": fichier,
                    "contenu": contenu[:500],
                    "analyse": resultat
                })

        # ===========================
        # 💾 Enregistrement des Résultats
        # ===========================
        output_domaine = os.path.join(OUTPUT_DIR, domaine)
        os.makedirs(output_domaine, exist_ok=True)
        output_file = os.path.join(output_domaine, f"{domaine}_analyses_gpt.json")

        with open(output_file, "w", encoding="utf-8") as f:
            json.dump(analyses, f, indent=4, ensure_ascii=False)

        print(f"✅ Analyse GPT terminée pour {domaine}. Résultats dans {output_file}")

# ===========================
# ▶️ Exécution
# ===========================
if __name__ == "__main__":
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    analyser_toutes_pages()
